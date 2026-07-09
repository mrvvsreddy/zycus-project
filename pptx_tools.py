"""
PPTX Tools — Reusable slide-building primitives for premium presentations.
--------------------------------------------------------------------------
Each function is a self-contained "tool" that adds one visual element
to a python-pptx slide. Compose them together in any slide builder.

Tools:
  Design tokens  — Colors, fonts, dimensions
  set_slide_bg   — Fill a slide with a solid dark background
  add_card       — Rounded rectangle card shape
  add_text       — Formatted text box with full font control
  add_accent_bar — Thin decorative accent line
  add_metric_card— KPI card: big number + label + accent stripe
  add_numbered_card — Numbered insight/recommendation row
  add_data_table — Styled table with header + alternating rows
  add_slide_header — Title + accent line combo (top of slide)
  embed_chart    — Place a matplotlib figure buffer onto a slide
  chart_donut    — RAG status donut chart → image buffer
  chart_hbar     — Horizontal bar chart → image buffer
  chart_grouped_vbar — Grouped vertical bar chart → image buffer
"""

import io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ═══════════════════════════════════════════════════════════════
#  DESIGN TOKENS
# ═══════════════════════════════════════════════════════════════
BG_DARK      = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD      = RGBColor(0x25, 0x25, 0x46)
ACCENT_CYAN  = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_AMBER = RGBColor(0xFF, 0xBE, 0x0B)
ACCENT_RED   = RGBColor(0xFF, 0x4D, 0x6A)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xB0, 0xB0, 0xC8)
MUTED        = RGBColor(0x70, 0x70, 0x90)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_TITLE = "Calibri"
FONT_BODY  = "Calibri"

RAG_COLORS = {"Red": ACCENT_RED, "Amber": ACCENT_AMBER, "Green": ACCENT_GREEN}
RAG_HEX    = {"Red": "#FF4D6A", "Amber": "#FFBE0B", "Green": "#00E696"}


# ═══════════════════════════════════════════════════════════════
#  PRIMITIVE TOOLS
# ═══════════════════════════════════════════════════════════════

def set_slide_bg(slide, color=None):
    """Tool: Fill the entire slide background with a solid color.
    Args:
        slide: pptx slide object
        color: RGBColor (defaults to BG_DARK)
    """
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or BG_DARK


def add_card(slide, left, top, width, height, fill_color=None):
    """Tool: Add a rounded-rectangle card shape.
    Args:
        slide: pptx slide object
        left, top, width, height: position/size in Inches or Emu
        fill_color: RGBColor (defaults to BG_CARD)
    Returns: the shape object
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text,
             font_size=14, color=None, bold=False,
             alignment=PP_ALIGN.LEFT, font_name=None):
    """Tool: Add a text box with full font control.
    Args:
        slide: pptx slide object
        left, top, width, height: position/size
        text: string content
        font_size: point size (default 14)
        color: RGBColor (defaults to WHITE)
        bold: boolean
        alignment: PP_ALIGN enum
        font_name: string (defaults to FONT_BODY)
    Returns: the textbox shape
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color or WHITE
    p.font.bold = bold
    p.font.name = font_name or FONT_BODY
    p.alignment = alignment
    return txBox


def add_accent_bar(slide, left, top, width, color=None):
    """Tool: Add a thin horizontal accent line (3pt tall).
    Args:
        slide: pptx slide object
        left, top, width: position/size
        color: RGBColor (defaults to ACCENT_CYAN)
    """
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = color or ACCENT_CYAN
    line.line.fill.background()
    return line


def add_vertical_bar(slide, left, top, height, color=None):
    """Tool: Add a thin vertical accent bar (3pt wide).
    Args:
        slide: pptx slide object
        left, top, height: position/size
        color: RGBColor (defaults to ACCENT_CYAN)
    """
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Pt(4), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color or ACCENT_CYAN
    bar.line.fill.background()
    return bar


# ═══════════════════════════════════════════════════════════════
#  COMPOSITE TOOLS (built from primitives)
# ═══════════════════════════════════════════════════════════════

def add_slide_header(slide, title, accent_color=None):
    """Tool: Add a slide title + accent underline at top-left.
    Args:
        slide: pptx slide object
        title: string
        accent_color: RGBColor for the underline
    """
    add_text(slide, Inches(0.7), Inches(0.4), Inches(10), Inches(0.6),
             title, font_size=28, bold=True, color=WHITE, font_name=FONT_TITLE)
    add_accent_bar(slide, Inches(0.7), Inches(1.05), Inches(2),
                   accent_color or ACCENT_CYAN)


def add_metric_card(slide, left, top, value, label,
                    accent_color=None, card_w=None, card_h=None):
    """Tool: Add a KPI metric card — big number on top, label below.
    Args:
        slide: pptx slide object
        left, top: position
        value: display value (string or number)
        label: description text below the number
        accent_color: color for the number + top stripe
        card_w, card_h: card dimensions (defaults to 2.8 × 1.6 inches)
    """
    cw = card_w or Inches(2.8)
    ch = card_h or Inches(1.6)
    ac = accent_color or ACCENT_CYAN

    add_card(slide, left, top, cw, ch)
    add_accent_bar(slide, left + Inches(0.25), top + Inches(0.15),
                   Inches(0.6), ac)
    add_text(slide, left + Inches(0.25), top + Inches(0.35),
             Inches(2.3), Inches(0.6),
             str(value), font_size=30, bold=True, color=ac)
    add_text(slide, left + Inches(0.25), top + Inches(0.95),
             Inches(2.3), Inches(0.4),
             label, font_size=11, color=LIGHT_GRAY)


def add_numbered_card(slide, left, top, width, number, text,
                      accent_color=None):
    """Tool: Add a numbered insight/recommendation row with accent bar.
    Args:
        slide: pptx slide object
        left, top, width: position/size
        number: integer to display
        text: content string
        accent_color: left-edge color
    """
    ac = accent_color or ACCENT_CYAN
    card_h = Inches(1.05)

    add_card(slide, left, top, width, card_h)
    add_vertical_bar(slide, left, top, card_h, ac)
    add_text(slide, left + Inches(0.4), top + Inches(0.18),
             Inches(0.5), Inches(0.5),
             str(number), font_size=22, bold=True, color=ac)
    add_text(slide, left + Inches(1.0), top + Inches(0.15),
             width - Inches(1.4), Inches(0.75),
             text, font_size=14, color=LIGHT_GRAY)


def add_data_table(slide, left, top, headers, rows, col_widths,
                   rag_col_index=None):
    """Tool: Add a styled data table with header row + alternating row colors.
    Args:
        slide: pptx slide object
        left, top: position
        headers: list of column header strings
        rows: list of dicts with keys matching columns
        col_widths: list of Inches values for each column
        rag_col_index: optional index of the column to color-code by RAG
    """
    total_w = sum(w for w in col_widths)
    row_h = Inches(0.6)

    # Header background
    add_card(slide, left, top, total_w, Inches(0.55))
    x = left
    for i, h in enumerate(headers):
        add_text(slide, x, top + Inches(0.08), col_widths[i], Inches(0.4),
                 h, font_size=12, bold=True, color=ACCENT_CYAN)
        x += col_widths[i]

    # Data rows
    for j, row_data in enumerate(rows):
        y = top + Inches(0.65) + row_h * j
        row_bg = BG_CARD if j % 2 == 0 else BG_DARK
        add_card(slide, left, y, total_w, row_h, row_bg)

        x = left
        for i, (key, val) in enumerate(row_data.items()):
            # Color-code RAG column
            txt_color = LIGHT_GRAY
            if rag_col_index is not None and i == rag_col_index:
                txt_color = RAG_COLORS.get(str(val), LIGHT_GRAY)
            add_text(slide, x + Inches(0.15), y + Inches(0.12),
                     col_widths[i] - Inches(0.15), Inches(0.4),
                     str(val), font_size=12, color=txt_color)
            x += col_widths[i]


def add_rag_badge(slide, left, top, status):
    """Tool: Add a small color-coded RAG status badge (circle).
    Args:
        slide: pptx slide object
        left, top: position
        status: 'Red', 'Amber', or 'Green'
    """
    color = RAG_COLORS.get(status, LIGHT_GRAY)
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()


# ═══════════════════════════════════════════════════════════════
#  CHART TOOLS (matplotlib → image buffer)
# ═══════════════════════════════════════════════════════════════

def _fig_to_buf(fig):
    """Convert a matplotlib figure to an in-memory PNG buffer."""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight",
                transparent=True, facecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf


def chart_donut(data_dict, colors_dict, center_text=None):
    """Tool: Generate a donut chart as a PNG image buffer.
    Args:
        data_dict: {"Green": 1, "Red": 1, ...} — label→count
        colors_dict: {"Green": "#00E696", ...} — label→hex color
        center_text: optional text to display in the center
    Returns: BytesIO buffer containing the PNG
    """
    labels, sizes, colors = [], [], []
    for label, count in data_dict.items():
        if count > 0:
            labels.append(label)
            sizes.append(count)
            colors.append(colors_dict.get(label, "#B0B0C8"))

    if not sizes:
        sizes, labels, colors = [1], ["N/A"], ["#555555"]

    fig, ax = plt.subplots(figsize=(4, 4))
    wedges, _, autotexts = ax.pie(
        sizes, labels=None, colors=colors, autopct="%1.0f%%",
        startangle=90, pctdistance=0.78,
        wedgeprops=dict(width=0.35, edgecolor="#1A1A2E", linewidth=2.5))

    for t in autotexts:
        t.set_color("white")
        t.set_fontsize(14)
        t.set_fontweight("bold")

    if center_text:
        ax.text(0, 0, center_text, ha="center", va="center",
                fontsize=12, color="white", fontweight="bold")

    ax.legend(labels, loc="center", fontsize=11,
              labelcolor="white", facecolor="none", edgecolor="none")
    ax.set_facecolor("none")
    fig.patch.set_alpha(0)
    return _fig_to_buf(fig)


def chart_hbar(names, values, colors, xlabel="", value_fmt="{:.1f}%",
              figsize=None):
    """Tool: Generate a horizontal bar chart as a PNG image buffer.
    Args:
        names: list of category names
        values: list of numeric values
        colors: list of hex color strings (one per bar)
        xlabel: x-axis label
        value_fmt: format string for value annotations
        figsize: (width, height) tuple in inches for the figure
    Returns: BytesIO buffer containing the PNG
    """
    if figsize is None:
        figsize = (5.5, max(2.5, len(names) * 1.1))
    fig, ax = plt.subplots(figsize=figsize)

    # Draw grid FIRST so bars render on top
    ax.grid(axis="x", color="#333355", linewidth=0.4, alpha=0.3, zorder=0)

    bars = ax.barh(names, values, color=colors, height=0.45,
                   edgecolor="none", linewidth=0, zorder=3)

    ax.set_xlim(0, max(values + [10]) * 1.3)
    ax.set_xlabel(xlabel, color="white", fontsize=11, labelpad=8)
    ax.invert_yaxis()
    ax.tick_params(colors="white", labelsize=12, length=0)
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.xaxis.set_major_formatter(mticker.FormatStrFormatter("%.0f%%"))

    # Value annotations to the right of each bar
    max_val = max(values + [10])
    for bar, val in zip(bars, values):
        label_x = bar.get_width() + max_val * 0.03
        ax.text(label_x, bar.get_y() + bar.get_height() / 2,
                value_fmt.format(val), va="center", color="white",
                fontsize=12, fontweight="bold", zorder=4)

    fig.patch.set_alpha(0)
    ax.set_facecolor("none")
    fig.tight_layout(pad=1.0)
    return _fig_to_buf(fig)


def chart_grouped_vbar(names, group_data, group_colors, ylabel="Count",
                      figsize=None):
    """Tool: Generate a grouped vertical bar chart as a PNG image buffer.
    Args:
        names: list of x-axis category names
        group_data: dict of {group_label: [values]} — one list per group
        group_colors: dict of {group_label: hex_color}
        ylabel: y-axis label
        figsize: (width, height) tuple in inches for the figure
    Returns: BytesIO buffer containing the PNG
    """
    import numpy as np

    n_groups = len(group_data)
    n_cats = len(names)
    w = 0.6 / n_groups
    x = np.arange(n_cats, dtype=float)

    if figsize is None:
        figsize = (5.5, 3.5)
    fig, ax = plt.subplots(figsize=figsize)

    # Draw grid FIRST (behind bars)
    ax.grid(axis="y", color="#333355", linewidth=0.4, alpha=0.3, zorder=0)

    # Force integer y-axis ticks
    max_val = max(max(v) for v in group_data.values()) if group_data else 1
    y_offset = max(max_val * 0.02, 0.2)

    for i, (label, vals) in enumerate(group_data.items()):
        offset = (i - n_groups / 2 + 0.5) * w
        positions = x + offset
        ax.bar(positions, vals, w, label=label,
               color=group_colors.get(label, "#888"),
               edgecolor="none", linewidth=0, zorder=3)
        # Value labels on top of bars
        for pos, val in zip(positions, vals):
            if val > 0:
                ax.text(pos, val + y_offset, str(int(val)),
                        ha="center", va="bottom", color="white",
                        fontsize=11, fontweight="bold", zorder=4)

    ax.set_xticks(x)
    ax.set_xticklabels(names, color="white", fontsize=12)
    ax.tick_params(axis="y", colors="white", labelsize=10, length=0)
    ax.tick_params(axis="x", length=0)
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.set_ylabel(ylabel, color="white", fontsize=10, labelpad=8)

    ax.set_ylim(0, max_val * 1.3)
    ax.yaxis.set_major_locator(mticker.MaxNLocator(integer=True))

    ax.legend(fontsize=10, labelcolor="white", facecolor="none",
              edgecolor="none", loc="upper right")
    fig.patch.set_alpha(0)
    ax.set_facecolor("none")
    fig.tight_layout(pad=1.0)
    return _fig_to_buf(fig)


def embed_chart(slide, chart_buf, left, top, width, height=None):
    """Tool: Place a chart image buffer onto a slide.
    Args:
        slide: pptx slide object
        chart_buf: BytesIO buffer from a chart_* function
        left, top, width: position/size
        height: optional (auto-scales if None)
    """
    if height:
        slide.shapes.add_picture(chart_buf, left, top, width, height)
    else:
        slide.shapes.add_picture(chart_buf, left, top, width=width)
