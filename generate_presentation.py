"""
Phase 3 – Monthly Synthesis Presentation Generator
---------------------------------------------------
Reads JSON reports from reports/, aggregates portfolio data,
calls the LLM for executive insights, and builds a premium
dark-themed PowerPoint using composable tools from pptx_tools.py.
"""

import os
import re
import json
import glob
from datetime import datetime

import requests
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches

from pptx_tools import (
    # Design tokens
    BG_DARK, BG_CARD, ACCENT_CYAN, ACCENT_GREEN, ACCENT_AMBER, ACCENT_RED,
    WHITE, LIGHT_GRAY, MUTED, SLIDE_W, SLIDE_H, FONT_TITLE, RAG_COLORS, RAG_HEX,
    # Primitive tools
    set_slide_bg, add_card, add_text, add_accent_bar, add_vertical_bar,
    # Composite tools
    add_slide_header, add_metric_card, add_numbered_card,
    add_data_table, add_rag_badge,
    # Chart tools
    chart_donut, chart_hbar, chart_grouped_vbar, embed_chart,
)

load_dotenv()
OPENROUTER_URL = "https://openrouter.ai/api/v1/chat/completions"
OPENROUTER_MODEL = os.environ.get("OPENROUTER_MODEL")


# ═══════════════════════════════════════════════════════════════
#  DATA LOADING & AGGREGATION
# ═══════════════════════════════════════════════════════════════

def load_reports(reports_dir="reports"):
    """Load all JSON project reports from the specified directory."""
    reports = []
    for f in sorted(glob.glob(os.path.join(reports_dir, "*.json"))):
        with open(f, "r") as fh:
            reports.append(json.load(fh))
    return reports


def aggregate_data(reports):
    """Aggregate metrics across all projects for portfolio-level analysis."""
    if not reports:
        return {}
    rag_counts = {"Red": 0, "Amber": 0, "Green": 0}
    total_tasks = total_active = total_overdue = total_blockers = total_crit = 0
    projects_risk = []

    for r in reports:
        rag = r.get("rag_status", "Green")
        if rag in rag_counts:
            rag_counts[rag] += 1
        facts = r.get("facts", {})
        total_tasks += facts.get("total_tasks", 0)
        total_active += facts.get("active_tasks", 0)
        total_overdue += facts.get("overdue_count", 0)
        total_blockers += facts.get("blocker_count", 0)
        total_crit += facts.get("critical_blocked_count", 0)
        projects_risk.append({
            "name": r.get("project_name", "Unknown"),
            "rag": rag,
            "pct_overdue": facts.get("pct_overdue") or 0.0,
            "blockers": facts.get("blocker_count", 0),
            "critical_blocked": facts.get("critical_blocked_count", 0),
            "active_tasks": facts.get("active_tasks", 0),
        })
    projects_risk.sort(
        key=lambda x: (x["critical_blocked"], x["pct_overdue"]), reverse=True)

    return {
        "total_projects": len(reports),
        "rag_counts": rag_counts,
        "total_tasks": total_tasks,
        "total_active_tasks": total_active,
        "total_overdue": total_overdue,
        "total_blockers": total_blockers,
        "total_critical_blocked": total_crit,
        "overall_pct_overdue": (
            round(100 * total_overdue / total_active, 1) if total_active else 0.0),
        "projects_risk": projects_risk,
    }


# ═══════════════════════════════════════════════════════════════
#  LLM SYNTHESIS
# ═══════════════════════════════════════════════════════════════

FALLBACK_INSIGHTS = [
    "Portfolio health is mixed — schedule slippage is concentrated in a subset of projects.",
    "Blocker count remains elevated; root-cause analysis is needed.",
    "Projects with Green status show strong execution discipline.",
]
FALLBACK_RECS = [
    "Escalate Red projects to steering committee this week.",
    "Conduct blocker-resolution workshops for affected teams.",
    "Institutionalise weekly RAG reviews to catch drift early.",
]


def _parse_llm_bullets(content):
    """Robustly parse LLM output into insights + recommendations lists.
    Handles: - dashes, * asterisks, 1. numbered, **bold** prefixes."""
    insights, recs = [], []
    section = None
    for line in content.split("\n"):
        stripped = line.strip()
        if not stripped:
            continue
        upper = stripped.upper()
        if "INSIGHT" in upper and (":" in stripped or stripped.endswith("S")):
            section = "i"
            continue
        if "RECOMMEND" in upper and (":" in stripped or stripped.endswith("S")):
            section = "r"
            continue
        # Strip leading bullet markers: -, *, 1., 1), •, **
        text = re.sub(r'^[\-\*•]+\s*', '', stripped)
        text = re.sub(r'^\d+[\.\)]\s*', '', text)
        text = re.sub(r'^\*\*(.+?)\*\*:?\s*', r'\1: ', text)
        text = text.strip()
        if not text:
            continue
        if section == "i":
            insights.append(text)
        elif section == "r":
            recs.append(text)
    return insights, recs


def generate_insights_with_llm(agg_data):
    """Call OpenRouter LLM and parse into (insights_list, recs_list)."""
    api_key = os.environ.get("OPENROUTER_API_KEY")
    if not api_key:
        return FALLBACK_INSIGHTS, FALLBACK_RECS

    prompt = f"""You are a VP of Delivery analyzing a portfolio of projects.

Aggregated data:
{json.dumps({k: v for k, v in agg_data.items() if k != 'projects_risk'}, indent=2)}

Top risk projects:
{json.dumps(agg_data['projects_risk'][:3], indent=2)}

Write EXACTLY 3 bullet points of Executive Insights (trends & risks).
Then EXACTLY 3 bullet points of Recommendations (actionable next steps).

Use this EXACT format:
INSIGHTS:
- First insight here
- Second insight here
- Third insight here
RECOMMENDATIONS:
- First recommendation here
- Second recommendation here
- Third recommendation here
"""
    try:
        resp = requests.post(
            OPENROUTER_URL,
            headers={"Authorization": f"Bearer {api_key}",
                     "Content-Type": "application/json"},
            json={"model": OPENROUTER_MODEL or "openai/gpt-4o-mini",
                  "messages": [{"role": "user", "content": prompt}],
                  "max_tokens": 600},
            timeout=30,
        )
        resp.raise_for_status()
        content = resp.json()["choices"][0]["message"]["content"].strip()
        insights, recs = _parse_llm_bullets(content)
        return (insights if len(insights) >= 2 else FALLBACK_INSIGHTS,
                recs if len(recs) >= 2 else FALLBACK_RECS)
    except Exception:
        return FALLBACK_INSIGHTS, FALLBACK_RECS


# ═══════════════════════════════════════════════════════════════
#  SLIDE BUILDERS (each uses tools from pptx_tools)
# ═══════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1: Title slide with accent bar and date."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Left accent strip
    add_card(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, ACCENT_CYAN)

    add_text(slide, Inches(1.2), Inches(2.0), Inches(10), Inches(1.2),
             "Monthly Portfolio Synthesis",
             font_size=44, bold=True, color=WHITE, font_name=FONT_TITLE)
    add_accent_bar(slide, Inches(1.2), Inches(3.3), Inches(3), ACCENT_CYAN)
    add_text(slide, Inches(1.2), Inches(3.65), Inches(8), Inches(0.5),
             f"Executive Summary  ·  {datetime.now().strftime('%B %Y')}",
             font_size=18, color=LIGHT_GRAY)
    add_text(slide, Inches(1.2), Inches(5.5), Inches(6), Inches(0.4),
             "Auto-generated by Project Health Agent",
             font_size=11, color=MUTED)


def slide_exec_summary(prs, agg):
    """Slide 2: KPI metric cards + RAG donut + RAG breakdown bars."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_header(slide, "Executive Summary")

    # Row of 4 metric cards
    add_metric_card(slide, Inches(0.7), Inches(1.4),
                    agg["total_projects"], "Active Projects", ACCENT_CYAN)
    add_metric_card(slide, Inches(3.8), Inches(1.4),
                    f"{agg['total_active_tasks']:,}", "Total Tasks", ACCENT_CYAN)
    add_metric_card(slide, Inches(6.9), Inches(1.4),
                    f"{agg['overall_pct_overdue']}%", "Overdue Rate", ACCENT_AMBER)
    add_metric_card(slide, Inches(10.0), Inches(1.4),
                    agg["total_blockers"], "Total Blockers", ACCENT_RED)

    # Donut chart
    donut_buf = chart_donut(agg["rag_counts"], RAG_HEX)
    embed_chart(slide, donut_buf, Inches(0.8), Inches(3.4),
                Inches(4.2), Inches(4.2))

    # RAG breakdown cards (right side)
    y = Inches(3.7)
    for status, color in [("Green", ACCENT_GREEN),
                          ("Amber", ACCENT_AMBER),
                          ("Red", ACCENT_RED)]:
        cnt = agg["rag_counts"].get(status, 0)
        add_card(slide, Inches(6.0), y, Inches(6.5), Inches(0.9))
        add_vertical_bar(slide, Inches(6.0), y, Inches(0.9), color)
        add_text(slide, Inches(6.5), y + Inches(0.12),
                 Inches(1.2), Inches(0.5),
                 str(cnt), font_size=28, bold=True, color=color)
        add_text(slide, Inches(7.8), y + Inches(0.2),
                 Inches(4), Inches(0.5),
                 f"{status} Projects", font_size=16, color=LIGHT_GRAY)
        y += Inches(1.1)


def slide_overdue_trends(prs, agg):
    """Slide 3: Overdue bar chart — sized to fit the slide perfectly."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_header(slide, "Portfolio Trends: Overdue Tasks")

    projects = agg["projects_risk"]
    names = [p["name"] for p in projects]
    pct_vals = [p["pct_overdue"] for p in projects]
    bar_colors = [RAG_HEX[p["rag"]] for p in projects]

    # ── Calculate available chart area on the slide ──
    # Slide: 13.333 × 7.5 inches
    # Margins: 0.7 left/right, header ~1.3 top, 0.5 bottom
    margin_lr = 0.7          # left & right margin
    header_h = 1.4           # header + accent bar height
    bottom_pad = 0.5         # bottom breathing room
    chart_left = margin_lr
    chart_top = header_h
    chart_w = 13.333 - 2 * margin_lr   # ~11.93 inches
    chart_h = 7.5 - header_h - bottom_pad  # ~5.6 inches

    # Create figure at matching aspect ratio (matplotlib inches)
    fig_scale = 1.0  # 1:1 mapping — matplotlib inch == slide inch
    fig_w = chart_w * fig_scale
    fig_h = chart_h * fig_scale

    hbar_buf = chart_hbar(names, pct_vals, bar_colors,
                          xlabel="% Overdue Tasks",
                          figsize=(fig_w, fig_h))
    embed_chart(slide, hbar_buf,
                Inches(chart_left), Inches(chart_top),
                Inches(chart_w), Inches(chart_h))


def slide_blocker_trends(prs, agg):
    """Slide 4: Blockers grouped bar chart — sized to fit the slide perfectly."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_header(slide, "Portfolio Trends: Blockers")

    projects = agg["projects_risk"]
    names = [p["name"] for p in projects]

    # ── Calculate available chart area on the slide ──
    margin_lr = 0.7
    header_h = 1.4
    bottom_pad = 0.9         # extra room for summary text
    chart_left = margin_lr
    chart_top = header_h
    chart_w = 13.333 - 2 * margin_lr   # ~11.93 inches
    chart_h = 7.5 - header_h - bottom_pad  # ~5.2 inches

    fig_w = chart_w
    fig_h = chart_h

    # Blockers grouped vertical bar chart
    blocker_vals = [p["blockers"] for p in projects]
    crit_vals = [p["critical_blocked"] for p in projects]
    vbar_buf = chart_grouped_vbar(
        names,
        {"Blockers": blocker_vals, "Critical Blocked": crit_vals},
        {"Blockers": "#00D2FF", "Critical Blocked": "#FF4D6A"},
        figsize=(fig_w, fig_h))
    embed_chart(slide, vbar_buf,
                Inches(chart_left), Inches(chart_top),
                Inches(chart_w), Inches(chart_h))

    # Summary text below the chart
    text_top = chart_top + chart_h + 0.1
    add_text(slide, Inches(margin_lr), Inches(text_top),
             Inches(chart_w), Inches(0.5),
             f"Portfolio-wide: {agg['total_blockers']} blockers, "
             f"{agg['total_critical_blocked']} critical blocked",
             font_size=12, color=MUTED)


def slide_risks(prs, agg):
    """Slide 4: Risk table with RAG color coding per project."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_header(slide, "Emerging Risks", ACCENT_RED)

    headers = ["Project", "RAG", "Overdue %", "Blockers",
               "Critical", "Active Tasks"]
    col_widths = [Inches(2.6), Inches(1.8), Inches(1.8),
                  Inches(1.6), Inches(1.5), Inches(2.2)]

    rows = []
    for p in agg["projects_risk"][:6]:
        rows.append({
            "Project": p["name"],
            "RAG": p["rag"],
            "Overdue %": f"{p['pct_overdue']:.1f}%",
            "Blockers": str(p["blockers"]),
            "Critical": str(p["critical_blocked"]),
            "Active Tasks": str(p["active_tasks"]),
        })

    add_data_table(slide, Inches(0.7), Inches(1.5),
                   headers, rows, col_widths, rag_col_index=1)

    # Add RAG badge dots next to project names
    for j, p in enumerate(agg["projects_risk"][:6]):
        y = Inches(1.5) + Inches(0.65) + Inches(0.6) * j
        add_rag_badge(slide, Inches(0.82), y + Inches(0.2), p["rag"])


def slide_insights(prs, insights_list):
    """Slide 5: Numbered executive insight cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_header(slide, "Executive Insights")

    for i, item in enumerate(insights_list[:4]):
        y = Inches(1.5) + Inches(1.3) * i
        add_numbered_card(slide, Inches(0.7), y, Inches(11.8),
                          i + 1, item, ACCENT_CYAN)


def slide_recommendations(prs, recs_list):
    """Slide 6: Numbered recommendation cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_header(slide, "Recommendations", ACCENT_GREEN)

    for i, item in enumerate(recs_list[:4]):
        y = Inches(1.5) + Inches(1.3) * i
        add_numbered_card(slide, Inches(0.7), y, Inches(11.8),
                          i + 1, item, ACCENT_GREEN)


def slide_closing(prs):
    """Slide 7: Thank-you / closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_card(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, ACCENT_CYAN)
    add_text(slide, Inches(1.2), Inches(2.5), Inches(10), Inches(1),
             "Thank You", font_size=44, bold=True, color=WHITE,
             font_name=FONT_TITLE)
    add_accent_bar(slide, Inches(1.2), Inches(3.6), Inches(2.5), ACCENT_CYAN)
    add_text(slide, Inches(1.2), Inches(4.0), Inches(8), Inches(0.5),
             f"Report generated {datetime.now().strftime('%B %d, %Y')}  ·  "
             f"Project Health Agent", font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
#  ORCHESTRATOR
# ═══════════════════════════════════════════════════════════════

def create_presentation(agg, insights, recs, out_file):
    """Compose all slides into a single presentation and save."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_exec_summary(prs, agg)
    slide_overdue_trends(prs, agg)
    slide_blocker_trends(prs, agg)
    slide_risks(prs, agg)
    slide_insights(prs, insights)
    slide_recommendations(prs, recs)
    slide_closing(prs)

    os.makedirs(os.path.dirname(out_file), exist_ok=True)
    prs.save(out_file)
    print(f"✅ Presentation saved to {out_file}")


def main():
    print("📂 Loading project reports...")
    reports = load_reports()
    if not reports:
        print("❌ No reports found in 'reports/'. Run the health agent first.")
        return

    print("📊 Aggregating portfolio data...")
    agg = aggregate_data(reports)

    print("🤖 Generating insights with LLM...")
    insights, recs = generate_insights_with_llm(agg)

    print("🎨 Building presentation...")
    out_file = (f"presentations/"
                f"Monthly_Synthesis_{datetime.now().strftime('%Y-%m-%d')}.pptx")
    create_presentation(agg, insights, recs, out_file)


if __name__ == "__main__":
    main()
